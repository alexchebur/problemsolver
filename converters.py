import io
import re
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
import PyPDF2
import pandas as pd
from typing import Optional

class WordToMarkdown:
    """Converter for Word documents to Markdown"""
    
    def convert(self, file_content: bytes) -> str:
        """Convert Word document content to Markdown"""
        doc = Document(io.BytesIO(file_content))
        markdown_content = []
        
        for paragraph in doc.paragraphs:
            text = paragraph.text.strip()
            if not text:
                markdown_content.append("")
                continue
            
            style_name = paragraph.style.name.lower()
            if 'heading' in style_name:
                if 'heading 1' in style_name:
                    markdown_content.append(f"# {text}")
                elif 'heading 2' in style_name:
                    markdown_content.append(f"## {text}")
                elif 'heading 3' in style_name:
                    markdown_content.append(f"### {text}")
                elif 'heading 4' in style_name:
                    markdown_content.append(f"#### {text}")
                elif 'heading 5' in style_name:
                    markdown_content.append(f"##### {text}")
                else:
                    markdown_content.append(f"###### {text}")
            else:
                formatted_text = self._format_text_runs(paragraph)
                markdown_content.append(formatted_text)
        
        for table in doc.tables:
            table_md = self._convert_table(table)
            markdown_content.append(table_md)
        
        return "\n\n".join(markdown_content)
    
    def _format_text_runs(self, paragraph) -> str:
        formatted_text = ""
        for run in paragraph.runs:
            text = run.text
            if run.bold and run.italic:
                text = f"***{text}***"
            elif run.bold:
                text = f"**{text}**"
            elif run.italic:
                text = f"*{text}*"
            formatted_text += text
        return formatted_text
    
    def _convert_table(self, table) -> str:
        markdown_table = []
        table_data = []
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                cell_text = cell.text.strip().replace('\n', ' ')
                row_data.append(cell_text)
            table_data.append(row_data)
        
        if not table_data:
            return ""
        
        header_row = "| " + " | ".join(table_data[0]) + " |"
        markdown_table.append(header_row)
        separator = "| " + " | ".join(["---"] * len(table_data[0])) + " |"
        markdown_table.append(separator)
        
        for row in table_data[1:]:
            data_row = "| " + " | ".join(row) + " |"
            markdown_table.append(data_row)
        
        return "\n".join(markdown_table)

class ExcelToMarkdown:
    """Converter for Excel files to Markdown"""
    
    def convert(self, file_content: bytes, for_analysis: bool = False) -> str:
        """Convert Excel file content to Markdown"""
        excel_io = io.BytesIO(file_content)
        workbook = load_workbook(excel_io, data_only=True)  # Используем data_only для получения значений вместо формул
        markdown_content = []
        
        for sheet_name in workbook.sheetnames:
            worksheet = workbook[sheet_name]
            
            # Для анализа временных рядов не включаем названия листов
            if not for_analysis:
                markdown_content.append(f"# {sheet_name}")
                markdown_content.append("")
            
            data = []
            for row in worksheet.iter_rows(values_only=True):
                if any(cell is not None and str(cell).strip() for cell in row):
                    clean_row = []
                    for cell in row:
                        if cell is None:
                            clean_row.append("")
                        else:
                            # Для числовых значений форматируем с ограниченной точностью
                            if isinstance(cell, (int, float)):
                                cell_value = f"{cell:.4f}"  # Ограничиваем до 4 знаков после запятой
                            else:
                                cell_value = str(cell).strip()
                            clean_row.append(cell_value)
                    data.append(clean_row)
            
            if data:
                table_md = self._create_markdown_table(data)
                markdown_content.append(table_md)
            elif not for_analysis:
                markdown_content.append("*No data in this worksheet*")
            
            markdown_content.append("")
        
        return "\n".join(markdown_content)

class PowerPointToMarkdown:
    """Converter for PowerPoint presentations to Markdown"""
    
    def convert(self, file_content: bytes) -> str:
        ppt_io = io.BytesIO(file_content)
        presentation = Presentation(ppt_io)
        markdown_content = ["# PowerPoint Presentation", ""]
        
        for slide_num, slide in enumerate(presentation.slides, 1):
            markdown_content.append(f"## Slide {slide_num}")
            markdown_content.append("")
            
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_content = shape.text.strip()
                    if not slide_text and self._is_likely_title(shape):
                        slide_text.append(f"### {text_content}")
                    else:
                        paragraphs = text_content.split('\n')
                        for para in paragraphs:
                            para = para.strip()
                            if para:
                                if para.startswith(('•', '-', '*')) or para[0:2] in ['1.', '2.', '3.', '4.', '5.']:
                                    slide_text.append(f"- {para.lstrip('•-* ')}")
                                else:
                                    slide_text.append(para)
                
                elif hasattr(shape, "table"):
                    table_md = self._convert_ppt_table(shape.table)
                    if table_md:
                        slide_text.append(table_md)
            
            if slide_text:
                markdown_content.extend(slide_text)
            else:
                markdown_content.append("*No text content in this slide*")
            
            markdown_content.append("")
        
        return "\n".join(markdown_content)
    
    def _is_likely_title(self, shape) -> bool:
        try:
            if hasattr(shape, 'top') and shape.top < 1000000:
                return True
        except:
            pass
        return False
    
    def _convert_ppt_table(self, table) -> str:
        markdown_table = []
        table_data = []
        
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                cell_text = cell.text.strip().replace('\n', ' ')
                row_data.append(cell_text)
            table_data.append(row_data)
        
        if not table_data:
            return ""
        
        if table_data:
            header_row = "| " + " | ".join(table_data[0]) + " |"
            markdown_table.append(header_row)
            separator = "| " + " | ".join(["---"] * len(table_data[0])) + " |"
            markdown_table.append(separator)
            
            for row in table_data[1:]:
                data_row = "| " + " | ".join(row) + " |"
                markdown_table.append(data_row)
        
        return "\n".join(markdown_table)

class PDFToMarkdown:
    """Converter for PDF documents to Markdown"""
    
    def convert(self, file_content: bytes) -> str:
        pdf_io = io.BytesIO(file_content)
        pdf_reader = PyPDF2.PdfReader(pdf_io)
        markdown_content = ["# PDF Document", ""]
        
        for page_num, page in enumerate(pdf_reader.pages, 1):
            markdown_content.append(f"## Page {page_num}")
            markdown_content.append("")
            
            try:
                page_text = page.extract_text()
                if page_text.strip():
                    cleaned_text = self._clean_pdf_text(page_text)
                    paragraphs = cleaned_text.split('\n\n')
                    
                    for paragraph in paragraphs:
                        paragraph = paragraph.strip()
                        if paragraph:
                            if self._is_likely_heading(paragraph):
                                markdown_content.append(f"### {paragraph}")
                            else:
                                markdown_content.append(paragraph)
                            markdown_content.append("")
                else:
                    markdown_content.append("*No text content found on this page*")
                    markdown_content.append("")
                    
            except Exception as e:
                markdown_content.append(f"*Error extracting text from page {page_num}: {str(e)}*")
                markdown_content.append("")
        
        return "\n".join(markdown_content)
    
    def _clean_pdf_text(self, text: str) -> str:
        if not text:
            return ""
        
        cleaned = ' '.join(text.split())
        cleaned = re.sub(r'([.!?])\s+([A-Z][a-z])', r'\1\n\n\2', cleaned)
        cleaned = re.sub(r'\s*([•·▪▫‣⁃])\s*', r'\n- ', cleaned)
        cleaned = re.sub(r'\s*(\d+\.)\s+', r'\n\1 ', cleaned)
        return cleaned
    
    def _is_likely_heading(self, text: str) -> bool:
        if len(text) < 100 and (
            text.isupper() or
            (len(text.split()) <= 8 and text.count('.') == 0)
        ):
            return True
        return False

class ConverterFactory:
    """Factory for file converters"""
    
    @staticmethod
    def get_converter(file_type: str):
        converters = {
            'docx': WordToMarkdown(),
            'xlsx': ExcelToMarkdown(),
            'pptx': PowerPointToMarkdown(),
            'pdf': PDFToMarkdown()
        }
        return converters.get(file_type)

def convert_uploaded_file_to_markdown(uploaded_file, for_analysis: bool = False) -> Optional[str]:
    """Parse uploaded file and convert to Markdown"""
    try:
        if uploaded_file is None:
            return None

        file_name = uploaded_file.name.lower()
        file_extension = file_name.split('.')[-1]
        
        converter = ConverterFactory.get_converter(file_extension)
        if not converter:
            return None

        file_content = uploaded_file.getvalue()
        
        # Для Excel используем специальный флаг
        if file_extension == 'xlsx' and for_analysis:
            return converter.convert(file_content, for_analysis=True)[:100000]  # Более строгое ограничение
        
        return converter.convert(file_content)[:300000]  # Limit to 300k characters

    except Exception as e:
        return None

def convert_excel_to_markdown_for_analysis(file_content: bytes, max_rows: int = 50) -> str:
    """Специальная конвертация Excel файла для анализа временных рядов"""
    try:
        excel_io = io.BytesIO(file_content)
        workbook = load_workbook(excel_io, data_only=True)
        markdown_content = []
        
        for sheet_name in workbook.sheetnames:
            worksheet = workbook[sheet_name]
            data = []
            headers = []
            
            # Получаем заголовки из первой строки
            if worksheet.max_row > 0:
                for cell in worksheet[1]:
                    headers.append(str(cell.value) if cell.value else f"Column{cell.column}")
            
            # Собираем данные (числовые и текстовые)
            for i, row in enumerate(worksheet.iter_rows(min_row=2, values_only=True), 2):
                if i > max_rows + 1:  # Ограничиваем количество строк
                    break
                
                row_data = []
                for j, cell in enumerate(row):
                    if cell is None:
                        row_data.append("")
                    elif isinstance(cell, (int, float)):
                        # Форматируем числа с 4 знаками после запятой
                        row_data.append(f"{cell:.4f}")
                    else:
                        # Сохраняем текстовые значения
                        row_data.append(str(cell).strip())
                data.append(row_data)
            
            if headers or data:
                markdown_content.append(f"## Лист: {sheet_name}")
                
                # Форматируем таблицу Markdown
                table = []
                if headers:
                    table.append("| " + " | ".join(headers) + " |")
                    table.append("| " + " | ".join(["---"] * len(headers)) + " |")
                
                for row in data:
                    table.append("| " + " | ".join(row) + " |")
                
                markdown_content.append("\n".join(table))
                markdown_content.append("")
        
        return "\n".join(markdown_content)[:200000]  # Ограничиваем размер
    
    except Exception as e:
        return f"Ошибка конвертации: {str(e)}"
